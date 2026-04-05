"""
Generate a professional PowerPoint slide deck from trend analysis data.

Creates a 7-slide .pptx report with charts embedded as images:
  1. Title
  2. Executive Summary
  3. Trending Topics (bar chart)
  4. Top Videos
  5. AI Tools in Spotlight (bar chart)
  6. Channel Activity
  7. Key Takeaways

Usage:
    python tools/create_report.py
    python tools/create_report.py .tmp/trends_2026-04-05.json
    python tools/create_report.py .tmp/trends_2026-04-05.json --output .tmp/report.pptx
"""

import argparse
import json
import os
import sys
from datetime import datetime, timezone
from glob import glob
from io import BytesIO

import matplotlib
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

matplotlib.use("Agg")  # Non-interactive backend

PROJECT_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

# ---------------------------------------------------------------------------
# Design constants
# ---------------------------------------------------------------------------
BG_COLOR = RGBColor(0x0F, 0x17, 0x2A)       # Dark navy
ACCENT = RGBColor(0x63, 0x66, 0xF1)          # Indigo
ACCENT2 = RGBColor(0x06, 0xB6, 0xD4)         # Cyan
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x94, 0xA3, 0xB8)
LIGHT_BG = RGBColor(0x1E, 0x29, 0x3B)        # Slightly lighter panel bg

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# matplotlib colors matching design
MPL_BG = "#0F172A"
MPL_PANEL = "#1E293B"
MPL_ACCENT = "#6366F1"
MPL_ACCENT2 = "#06B6D4"
MPL_WHITE = "#FFFFFF"
MPL_GRAY = "#94A3B8"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def add_blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]  # Blank layout
    return prs.slides.add_slide(blank_layout)


def fill_slide_bg(slide, color: RGBColor):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text: str, left, top, width, height,
                font_size=Pt(14), bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_accent_bar(slide, left, top, width=Inches(0.05), height=Inches(0.5),
                   color=ACCENT):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def chart_to_image(fig) -> BytesIO:
    buf = BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor=MPL_BG, edgecolor="none")
    buf.seek(0)
    plt.close(fig)
    return buf


def format_views(n: int) -> str:
    if n >= 1_000_000:
        return f"{n / 1_000_000:.1f}M"
    if n >= 1_000:
        return f"{n / 1_000:.0f}K"
    return str(n)


def add_divider_line(slide, top, color=ACCENT):
    line = slide.shapes.add_shape(1,
        Inches(0.5), top, Inches(12.33), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def build_title_slide(prs, trends: dict):
    slide = add_blank_slide(prs)
    fill_slide_bg(slide, BG_COLOR)

    date_from = trends["date_range"].get("from", "")
    date_to = trends["date_range"].get("to", "")
    total_videos = trends["summary"]["total_videos"]
    channels = trends["summary"]["channels_active"]

    # Large accent rectangle left edge
    add_accent_bar(slide, Inches(0), Inches(0), Inches(0.25), SLIDE_H, ACCENT)

    # Main title
    add_textbox(slide, "AI Trend Report",
                Inches(1), Inches(1.5), Inches(11), Inches(1.8),
                font_size=Pt(54), bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Subtitle line
    add_textbox(slide, f"{date_from}  →  {date_to}",
                Inches(1), Inches(3.4), Inches(8), Inches(0.6),
                font_size=Pt(22), bold=False, color=ACCENT2, align=PP_ALIGN.LEFT)

    # Stats row
    stats_text = f"{channels} channels  ·  {total_videos} videos analyzed"
    add_textbox(slide, stats_text,
                Inches(1), Inches(4.2), Inches(8), Inches(0.5),
                font_size=Pt(16), color=GRAY, align=PP_ALIGN.LEFT)

    # Bottom tagline
    add_textbox(slide, "What's trending in AI right now",
                Inches(1), Inches(5.8), Inches(10), Inches(0.5),
                font_size=Pt(13), color=GRAY, align=PP_ALIGN.LEFT)


def build_exec_summary_slide(prs, trends: dict):
    slide = add_blank_slide(prs)
    fill_slide_bg(slide, BG_COLOR)

    add_accent_bar(slide, Inches(0), Inches(0), Inches(0.25), SLIDE_H, ACCENT)
    add_textbox(slide, "Executive Summary",
                Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
                font_size=Pt(28), bold=True, color=WHITE)
    add_divider_line(slide, Inches(1.1))

    top_topics = trends["summary"]["top_topics"]
    topic_breakdown = trends["topic_breakdown"]
    total_videos = trends["summary"]["total_videos"]

    # 3 stat cards
    card_colors = [ACCENT, ACCENT2, RGBColor(0xF5, 0x9E, 0x0B)]
    card_data = []
    for topic in top_topics[:3]:
        data = topic_breakdown.get(topic, {})
        card_data.append({
            "topic": topic,
            "video_count": data.get("video_count", 0),
            "avg_views": data.get("avg_views", 0),
            "top_tools": data.get("top_tools", []),
        })

    card_left_starts = [Inches(0.8), Inches(4.8), Inches(8.8)]
    card_width = Inches(3.7)

    for i, (card, color) in enumerate(zip(card_data, card_colors)):
        left = card_left_starts[i]

        # Card background
        bg = slide.shapes.add_shape(1, left, Inches(1.4), card_width, Inches(4.8))
        bg.fill.solid()
        bg.fill.fore_color.rgb = LIGHT_BG
        bg.line.fill.background()

        # Top accent bar on card
        top_bar = slide.shapes.add_shape(1, left, Inches(1.4), card_width, Inches(0.1))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = color
        top_bar.line.fill.background()

        # Topic name
        add_textbox(slide, card["topic"],
                    left + Inches(0.2), Inches(1.6), card_width - Inches(0.4), Inches(0.8),
                    font_size=Pt(16), bold=True, color=WHITE)

        # Video count
        add_textbox(slide, f"{card['video_count']} videos",
                    left + Inches(0.2), Inches(2.5), card_width - Inches(0.4), Inches(0.5),
                    font_size=Pt(22), bold=True, color=color)

        add_textbox(slide, "in the last 14 days",
                    left + Inches(0.2), Inches(3.05), card_width - Inches(0.4), Inches(0.4),
                    font_size=Pt(11), color=GRAY)

        add_textbox(slide, f"Avg. {format_views(card['avg_views'])} views",
                    left + Inches(0.2), Inches(3.55), card_width - Inches(0.4), Inches(0.4),
                    font_size=Pt(13), color=WHITE)

        if card["top_tools"]:
            tools_str = "  ·  ".join(card["top_tools"])
            add_textbox(slide, f"Trending: {tools_str}",
                        left + Inches(0.2), Inches(4.1), card_width - Inches(0.4), Inches(0.8),
                        font_size=Pt(11), color=GRAY)

    # Footer stat
    add_textbox(slide, f"Based on {total_videos} videos across {trends['summary']['channels_active']} top AI channels",
                Inches(0.8), Inches(6.7), Inches(12), Inches(0.4),
                font_size=Pt(11), color=GRAY, align=PP_ALIGN.CENTER)


def build_trending_topics_slide(prs, trends: dict):
    slide = add_blank_slide(prs)
    fill_slide_bg(slide, BG_COLOR)

    add_accent_bar(slide, Inches(0), Inches(0), Inches(0.25), SLIDE_H, ACCENT)
    add_textbox(slide, "Trending Topics",
                Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
                font_size=Pt(28), bold=True, color=WHITE)
    add_divider_line(slide, Inches(1.1))

    topic_breakdown = {k: v for k, v in trends["topic_breakdown"].items() if k != "Other"}
    topics = list(topic_breakdown.keys())[:8]
    counts = [topic_breakdown[t]["video_count"] for t in topics]
    views = [topic_breakdown[t]["total_views"] for t in topics]

    # Sort by video count
    sorted_pairs = sorted(zip(topics, counts, views), key=lambda x: x[1])
    topics_s, counts_s, views_s = zip(*sorted_pairs) if sorted_pairs else ([], [], [])

    fig, ax = plt.subplots(figsize=(8.5, 4.2))
    fig.patch.set_facecolor(MPL_BG)
    ax.set_facecolor(MPL_PANEL)

    bars = ax.barh(topics_s, counts_s, color=MPL_ACCENT, height=0.6)
    ax.bar_label(bars, fmt="%d", padding=4, color=MPL_WHITE, fontsize=10)

    ax.set_xlabel("Number of Videos", color=MPL_GRAY, fontsize=10)
    ax.tick_params(colors=MPL_WHITE, labelsize=10)
    ax.spines[:].set_color(MPL_PANEL)
    ax.xaxis.label.set_color(MPL_GRAY)
    ax.set_title("Video Count by Topic (last 14 days)", color=MPL_WHITE, fontsize=12, pad=10)
    plt.tight_layout()

    img = chart_to_image(fig)
    slide.shapes.add_picture(img, Inches(0.7), Inches(1.25), Inches(8.7), Inches(5.8))

    # Side panel: top topics list
    for i, (topic, count) in enumerate(zip(reversed(list(topics_s)), reversed(list(counts_s)))):
        if i >= 5:
            break
        y = Inches(1.5) + i * Inches(0.85)
        add_textbox(slide, f"#{i+1}", Inches(9.7), y, Inches(0.4), Inches(0.4),
                    font_size=Pt(11), bold=True, color=ACCENT)
        add_textbox(slide, topic, Inches(10.1), y, Inches(3.0), Inches(0.4),
                    font_size=Pt(11), color=WHITE)
        add_textbox(slide, f"{count} videos", Inches(10.1), y + Inches(0.35),
                    Inches(2.5), Inches(0.35), font_size=Pt(10), color=GRAY)


def build_top_videos_slide(prs, trends: dict):
    slide = add_blank_slide(prs)
    fill_slide_bg(slide, BG_COLOR)

    add_accent_bar(slide, Inches(0), Inches(0), Inches(0.25), SLIDE_H, ACCENT)
    add_textbox(slide, "Top Videos This Period",
                Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
                font_size=Pt(28), bold=True, color=WHITE)
    add_divider_line(slide, Inches(1.1))

    top_videos = trends.get("top_videos", [])

    # Column headers
    headers = ["Title", "Channel", "Views", "Topic"]
    col_x = [Inches(0.8), Inches(6.5), Inches(9.0), Inches(10.6)]
    col_w = [Inches(5.5), Inches(2.3), Inches(1.4), Inches(2.5)]

    for i, (header, x) in enumerate(zip(headers, col_x)):
        add_textbox(slide, header, x, Inches(1.25), col_w[i], Inches(0.35),
                    font_size=Pt(11), bold=True, color=ACCENT)

    add_divider_line(slide, Inches(1.62), color=LIGHT_BG)

    for row_i, video in enumerate(top_videos[:5]):
        y = Inches(1.7) + row_i * Inches(0.95)
        row_color = LIGHT_BG if row_i % 2 == 0 else BG_COLOR

        # Row background
        bg = slide.shapes.add_shape(1, Inches(0.7), y - Inches(0.05),
                                     Inches(12.5), Inches(0.85))
        bg.fill.solid()
        bg.fill.fore_color.rgb = row_color
        bg.line.fill.background()

        title = video["title"]
        if len(title) > 60:
            title = title[:57] + "..."
        topic_label = video["topics"][0] if video.get("topics") else ""

        add_textbox(slide, title, col_x[0], y, col_w[0], Inches(0.8),
                    font_size=Pt(11), color=WHITE)
        add_textbox(slide, video["channel_name"], col_x[1], y, col_w[1], Inches(0.8),
                    font_size=Pt(11), color=ACCENT2)
        add_textbox(slide, format_views(video["view_count"]), col_x[2], y, col_w[2], Inches(0.8),
                    font_size=Pt(11), bold=True, color=WHITE)
        add_textbox(slide, topic_label, col_x[3], y, col_w[3], Inches(0.8),
                    font_size=Pt(10), color=GRAY)


def build_tools_spotlight_slide(prs, trends: dict):
    slide = add_blank_slide(prs)
    fill_slide_bg(slide, BG_COLOR)

    add_accent_bar(slide, Inches(0), Inches(0), Inches(0.25), SLIDE_H, ACCENT)
    add_textbox(slide, "AI Tools in Spotlight",
                Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
                font_size=Pt(28), bold=True, color=WHITE)
    add_divider_line(slide, Inches(1.1))

    tools = trends.get("tools_in_spotlight", [])[:10]
    if not tools:
        add_textbox(slide, "No tool mentions found in this period.",
                    Inches(1), Inches(2), Inches(10), Inches(1),
                    font_size=Pt(16), color=GRAY)
        return

    tool_names = [t["tool"] for t in tools]
    mention_counts = [t["mention_count"] for t in tools]
    channels_mentioning = [t["channels_mentioning"] for t in tools]

    # Sort ascending for horizontal bar
    sorted_triples = sorted(zip(tool_names, mention_counts, channels_mentioning), key=lambda x: x[1])
    names_s, counts_s, chans_s = zip(*sorted_triples)

    fig, ax = plt.subplots(figsize=(8.0, 4.5))
    fig.patch.set_facecolor(MPL_BG)
    ax.set_facecolor(MPL_PANEL)

    # Color bars by channel spread
    max_chans = max(chans_s) if max(chans_s) > 0 else 1
    bar_colors = [MPL_ACCENT2 if c / max_chans >= 0.6 else MPL_ACCENT for c in chans_s]

    bars = ax.barh(names_s, counts_s, color=bar_colors, height=0.6)
    ax.bar_label(bars, fmt="%d", padding=4, color=MPL_WHITE, fontsize=10)

    ax.set_xlabel("Total Mentions", color=MPL_GRAY, fontsize=10)
    ax.tick_params(colors=MPL_WHITE, labelsize=10)
    ax.spines[:].set_color(MPL_PANEL)
    ax.set_title("Tool Mentions Across All Channels", color=MPL_WHITE, fontsize=12, pad=10)

    legend_handles = [
        mpatches.Patch(color=MPL_ACCENT2, label="High channel spread (≥60%)"),
        mpatches.Patch(color=MPL_ACCENT, label="Lower channel spread"),
    ]
    ax.legend(handles=legend_handles, facecolor=MPL_BG, labelcolor=MPL_WHITE,
              fontsize=9, loc="lower right")
    plt.tight_layout()

    img = chart_to_image(fig)
    slide.shapes.add_picture(img, Inches(0.7), Inches(1.25), Inches(8.3), Inches(5.9))

    # Side: top 3 spotlight callouts
    for i, tool_data in enumerate(tools[:3]):
        y = Inches(1.5) + i * Inches(1.5)
        add_textbox(slide, tool_data["tool"],
                    Inches(9.3), y, Inches(3.8), Inches(0.5),
                    font_size=Pt(16), bold=True, color=WHITE)
        add_textbox(slide,
                    f"{tool_data['mention_count']} mentions · {tool_data['channels_mentioning']} channels",
                    Inches(9.3), y + Inches(0.5), Inches(3.8), Inches(0.4),
                    font_size=Pt(11), color=GRAY)


def build_channel_activity_slide(prs, trends: dict):
    slide = add_blank_slide(prs)
    fill_slide_bg(slide, BG_COLOR)

    add_accent_bar(slide, Inches(0), Inches(0), Inches(0.25), SLIDE_H, ACCENT)
    add_textbox(slide, "Channel Activity",
                Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
                font_size=Pt(28), bold=True, color=WHITE)
    add_divider_line(slide, Inches(1.1))

    channels = trends.get("channel_activity", [])

    headers = ["Channel", "Videos Posted", "Avg Views", "Top Topic"]
    col_x = [Inches(0.8), Inches(4.5), Inches(7.0), Inches(9.5)]
    col_w = [Inches(3.5), Inches(2.3), Inches(2.3), Inches(3.6)]

    for header, x, w in zip(headers, col_x, col_w):
        add_textbox(slide, header, x, Inches(1.25), w, Inches(0.35),
                    font_size=Pt(11), bold=True, color=ACCENT)

    add_divider_line(slide, Inches(1.62), color=LIGHT_BG)

    for row_i, ch in enumerate(channels[:8]):
        y = Inches(1.7) + row_i * Inches(0.65)
        row_color = LIGHT_BG if row_i % 2 == 0 else BG_COLOR

        bg = slide.shapes.add_shape(1, Inches(0.7), y - Inches(0.04),
                                     Inches(12.5), Inches(0.6))
        bg.fill.solid()
        bg.fill.fore_color.rgb = row_color
        bg.line.fill.background()

        add_textbox(slide, ch["channel_name"], col_x[0], y, col_w[0], Inches(0.55),
                    font_size=Pt(11), bold=True, color=WHITE)
        add_textbox(slide, str(ch["videos_posted"]), col_x[1], y, col_w[1], Inches(0.55),
                    font_size=Pt(11), color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, format_views(ch["avg_views"]), col_x[2], y, col_w[2], Inches(0.55),
                    font_size=Pt(11), color=ACCENT2, align=PP_ALIGN.CENTER)
        add_textbox(slide, ch.get("top_topic", ""), col_x[3], y, col_w[3], Inches(0.55),
                    font_size=Pt(10), color=GRAY)


def build_takeaways_slide(prs, trends: dict):
    slide = add_blank_slide(prs)
    fill_slide_bg(slide, BG_COLOR)

    add_accent_bar(slide, Inches(0), Inches(0), Inches(0.25), SLIDE_H, ACCENT)
    add_textbox(slide, "Key Takeaways",
                Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
                font_size=Pt(28), bold=True, color=WHITE)
    add_divider_line(slide, Inches(1.1))

    top_topics = trends["summary"]["top_topics"]
    topic_breakdown = trends["topic_breakdown"]
    tools = trends.get("tools_in_spotlight", [])
    total_videos = trends["summary"]["total_videos"]
    channels = trends["summary"]["channels_active"]

    # Synthesize takeaways from data
    takeaways = []

    if top_topics:
        t1 = top_topics[0]
        count1 = topic_breakdown.get(t1, {}).get("video_count", 0)
        pct = round(count1 / total_videos * 100) if total_videos else 0
        takeaways.append(
            f"{t1} dominates the conversation — {count1} of {total_videos} videos "
            f"({pct}%) cover this topic, signaling strong audience demand."
        )

    if len(top_topics) >= 2:
        t2 = top_topics[1]
        tools_t2 = topic_breakdown.get(t2, {}).get("top_tools", [])
        if tools_t2:
            takeaways.append(
                f"{t2} content is surging, led by {', '.join(tools_t2[:2])}. "
                f"Consider creating comparison or tutorial content in this space."
            )
        else:
            takeaways.append(
                f"{t2} is the second most covered topic — worth exploring for your next content piece."
            )

    if tools:
        top_tool = tools[0]["tool"]
        top_count = tools[0]["mention_count"]
        top_chans = tools[0]["channels_mentioning"]
        takeaways.append(
            f"{top_tool} is the most discussed tool — mentioned {top_count} times across "
            f"{top_chans} channels. High creator interest = high audience interest."
        )

    takeaways.append(
        f"Tracking {channels} channels gave you visibility into {total_videos} videos. "
        f"Add more channels or widen the date window to increase coverage."
    )

    icons = ["01", "02", "03", "04"]
    for i, (takeaway, icon) in enumerate(zip(takeaways, icons)):
        y = Inches(1.4) + i * Inches(1.3)

        # Number badge
        badge = slide.shapes.add_shape(1, Inches(0.8), y, Inches(0.5), Inches(0.5))
        badge.fill.solid()
        badge.fill.fore_color.rgb = ACCENT
        badge.line.fill.background()

        add_textbox(slide, icon, Inches(0.8), y, Inches(0.5), Inches(0.5),
                    font_size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        add_textbox(slide, takeaway,
                    Inches(1.6), y, Inches(11.4), Inches(1.1),
                    font_size=Pt(13), color=WHITE, wrap=True)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def find_latest_trends_file() -> str | None:
    pattern = os.path.join(PROJECT_ROOT, ".tmp", "trends_*.json")
    files = sorted(glob(pattern), reverse=True)
    return files[0] if files else None


def main():
    parser = argparse.ArgumentParser(description="Generate PowerPoint trend report")
    parser.add_argument(
        "input",
        nargs="?",
        default=None,
        help="Path to trends_*.json (default: latest in .tmp/)",
    )
    parser.add_argument(
        "--output",
        default=None,
        help="Output .pptx path (default: .tmp/ai_trend_report_YYYY-MM-DD.pptx)",
    )
    args = parser.parse_args()

    input_path = args.input
    if not input_path:
        input_path = find_latest_trends_file()
        if not input_path:
            print("ERROR: No trends_*.json found in .tmp/. Run analyze_trends.py first.")
            sys.exit(1)
        print(f"Using latest trends file: {input_path}")

    with open(input_path) as f:
        trends = json.load(f)

    date_str = datetime.now(timezone.utc).strftime("%Y-%m-%d")
    output_path = args.output or os.path.join(
        PROJECT_ROOT, ".tmp", f"ai_trend_report_{date_str}.pptx"
    )
    os.makedirs(os.path.dirname(output_path), exist_ok=True)

    print("Building slides...")
    prs = new_presentation()

    print("  Slide 1: Title")
    build_title_slide(prs, trends)

    print("  Slide 2: Executive Summary")
    build_exec_summary_slide(prs, trends)

    print("  Slide 3: Trending Topics")
    build_trending_topics_slide(prs, trends)

    print("  Slide 4: Top Videos")
    build_top_videos_slide(prs, trends)

    print("  Slide 5: AI Tools in Spotlight")
    build_tools_spotlight_slide(prs, trends)

    print("  Slide 6: Channel Activity")
    build_channel_activity_slide(prs, trends)

    print("  Slide 7: Key Takeaways")
    build_takeaways_slide(prs, trends)

    prs.save(output_path)
    print(f"\nReport saved to {output_path}")
    return output_path


if __name__ == "__main__":
    main()
